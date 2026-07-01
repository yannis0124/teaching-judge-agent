from __future__ import annotations

import re
import shutil
from html.parser import HTMLParser
from pathlib import Path
from typing import Sequence

from src.processors.srt_processor import classify_text
from src.schemas.evidence_schema import PptSlideEvidence


def process_ppt(ppt_path: Path | Sequence[Path] | None, image_dir: Path) -> dict:
    image_dir.mkdir(parents=True, exist_ok=True)
    if ppt_path is None:
        return {
            "available": False,
            "text_available": False,
            "images_available": False,
            "status": "缺少slides.pptx或slides.html，课件页码证据不足，需要人工复核。",
            "image_status": "未生成课件页面截图。",
            "slides": [],
        }

    paths = [ppt_path] if isinstance(ppt_path, Path) else list(ppt_path)
    paths = [path for path in paths if path.exists()]
    if not paths:
        return {
            "available": False,
            "text_available": False,
            "images_available": False,
            "status": "缺少slides.pptx或slides.html，课件页码证据不足，需要人工复核。",
            "image_status": "未生成课件页面截图。",
            "slides": [],
        }
    if len(paths) == 1:
        return _process_single_courseware(paths[0], image_dir)

    merged_slides: list[dict] = []
    image_statuses: list[str] = []
    page = 1
    for index, path in enumerate(paths, start=1):
        result = _process_single_courseware(path, image_dir / f"courseware_{index:03d}")
        image_statuses.append(f"{path.name}: {result.get('image_status', '')}")
        for slide in result.get("slides", []):
            merged_slide = dict(slide)
            merged_slide["page"] = page
            merged_slide["text"] = f"来源文件：{path.name}\n{merged_slide.get('text', '')}".strip()
            merged_slides.append(merged_slide)
            page += 1

    return {
        "available": True,
        "text_available": bool(merged_slides),
        "images_available": any(bool(slide.get("image_path")) for slide in merged_slides),
        "status": f"已合并提取{len(paths)}个课件文件，共{len(merged_slides)}页课件文本。",
        "image_status": "；".join(image_statuses),
        "slides": merged_slides,
    }


def _process_single_courseware(ppt_path: Path, image_dir: Path) -> dict:
    if ppt_path.suffix.lower() in {".html", ".htm"}:
        return _process_html_courseware(ppt_path)

    try:
        slides = _extract_slide_text(ppt_path)
    except Exception as exc:
        return {
            "available": True,
            "text_available": False,
            "images_available": False,
            "status": f"PPT file exists but could not be parsed: {ppt_path.name}; {exc}",
            "image_status": "PPT images were not generated because the PPT file could not be parsed.",
            "slides": [],
        }
    image_status = _try_export_slide_images(ppt_path, image_dir)
    images_available = image_status["available"]
    for slide in slides:
        image_path = image_dir / f"slide_{slide.page:03d}.png"
        if image_path.exists():
            slide.image_path = str(image_path)

    return {
        "available": True,
        "text_available": True,
        "images_available": images_available,
        "status": f"已提取{len(slides)}页PPT文本和页码。",
        "image_status": image_status["status"],
        "slides": [slide.to_dict() for slide in slides],
    }


def _process_html_courseware(html_path: Path) -> dict:
    slides = _extract_html_text(html_path)
    return {
        "available": True,
        "text_available": bool(slides),
        "images_available": False,
        "status": f"已提取{len(slides)}页HTML课件文本。",
        "image_status": "HTML课件未生成页面截图；仅使用可见文本作为课件证据。",
        "slides": [slide.to_dict() for slide in slides],
    }


def _extract_slide_text(ppt_path: Path) -> list[PptSlideEvidence]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少python-pptx依赖，请先安装requirements.txt。") from exc

    presentation = Presentation(str(ppt_path))
    slides: list[PptSlideEvidence] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                texts.append(text.strip())
        merged_text = "\n".join(part for part in texts if part)
        slides.append(
            PptSlideEvidence(
                page=index,
                text=merged_text,
                tags=classify_text(merged_text),
                image_path=None,
            )
        )
    return slides


def _extract_html_text(html_path: Path) -> list[PptSlideEvidence]:
    parser = _CoursewareHtmlParser()
    parser.feed(html_path.read_text(encoding="utf-8", errors="ignore"))
    texts = parser.slide_texts()
    slides: list[PptSlideEvidence] = []
    for index, text in enumerate(texts, start=1):
        slides.append(
            PptSlideEvidence(
                page=index,
                text=text,
                tags=classify_text(text),
                image_path=None,
            )
        )
    return slides


class _CoursewareHtmlParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._skip_depth = 0
        self._section_depth = 0
        self._current_section: list[str] = []
        self._sections: list[str] = []
        self._document: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style", "noscript"}:
            self._skip_depth += 1
            return
        if self._skip_depth:
            return
        if tag == "section":
            if self._section_depth == 0:
                self._current_section = []
            self._section_depth += 1
            return
        if tag in {"br", "p", "div", "li", "h1", "h2", "h3", "h4", "tr"}:
            self._append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript"}:
            self._skip_depth = max(0, self._skip_depth - 1)
            return
        if self._skip_depth:
            return
        if tag == "section" and self._section_depth:
            self._section_depth -= 1
            if self._section_depth == 0:
                text = _normalize_html_text(" ".join(self._current_section))
                if text:
                    self._sections.append(text)
            return
        if tag in {"p", "div", "li", "h1", "h2", "h3", "h4", "tr"}:
            self._append("\n")

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return
        self._append(data)

    def slide_texts(self) -> list[str]:
        if self._sections:
            return self._sections
        document_text = _normalize_html_text(" ".join(self._document))
        return _chunk_text(document_text)

    def _append(self, text: str) -> None:
        if self._section_depth:
            self._current_section.append(text)
        else:
            self._document.append(text)


def _normalize_html_text(text: str) -> str:
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r"\n\s*", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _chunk_text(text: str, limit: int = 3000) -> list[str]:
    if not text:
        return []
    paragraphs = [part.strip() for part in re.split(r"\n+", text) if part.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        next_text = f"{current}\n{paragraph}".strip() if current else paragraph
        if len(next_text) <= limit:
            current = next_text
            continue
        if current:
            chunks.append(current)
        current = paragraph
    if current:
        chunks.append(current)
    return chunks


def _try_export_slide_images(ppt_path: Path, image_dir: Path) -> dict:
    com_result = _export_with_powerpoint_com(ppt_path, image_dir)
    if com_result["available"]:
        return com_result
    return {
        "available": False,
        "status": "PPT页面截图未生成，视觉证据不足，需要人工复核。",
    }


def _export_with_powerpoint_com(ppt_path: Path, image_dir: Path) -> dict:
    if shutil.which("powershell") is None:
        return {"available": False, "status": "无法调用PowerPoint导出PPT截图。"}
    try:
        import win32com.client  # type: ignore
    except ImportError:
        return {"available": False, "status": "缺少pywin32，PPT页面截图未生成，视觉证据不足，需要人工复核。"}

    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(str(ppt_path.resolve()), WithWindow=False)
        for idx, slide in enumerate(presentation.Slides, start=1):
            out_path = image_dir / f"slide_{idx:03d}.png"
            slide.Export(str(out_path.resolve()), "PNG")
        return {"available": True, "status": "已导出PPT页面截图。"}
    except Exception:
        return {
            "available": False,
            "status": "PPT页面截图未生成，视觉证据不足，需要人工复核。",
        }
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass
